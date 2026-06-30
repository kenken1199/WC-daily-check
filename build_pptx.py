# -*- coding: utf-8 -*-
"""WC分析ツール 紹介資料 PowerPoint生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy
import os

# ===== カラー定義 =====
C_NAVY      = RGBColor(0x1F, 0x38, 0x64)   # ネイビー
C_BLUE      = RGBColor(0x44, 0x72, 0xC4)   # ブルー
C_LTBLUE    = RGBColor(0x5B, 0x9B, 0xD5)   # ライトブルー
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_BGWARM    = RGBColor(0xFA, 0xF6, 0xF0)   # 温かみある白
C_BGSUB     = RGBColor(0xEF, 0xF3, 0xFB)   # サブ背景
C_TEXT      = RGBColor(0x2C, 0x3E, 0x50)
C_SUBTEXT   = RGBColor(0x7D, 0x8A, 0x93)
C_NG_BG     = RGBColor(0xFF, 0xE4, 0xE1)
C_NG_TEXT   = RGBColor(0x9C, 0x00, 0x06)
C_WARN_BG   = RGBColor(0xFF, 0xF2, 0xCC)
C_WARN_TEXT = RGBColor(0x7F, 0x60, 0x00)
C_OK_BG     = RGBColor(0xD4, 0xED, 0xDA)
C_OK_TEXT   = RGBColor(0x15, 0x57, 0x24)
C_GREEN     = RGBColor(0x28, 0xA7, 0x45)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT    = "Meiryo"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank

# ===== ヘルパー =====

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=Pt(14), bold=False, color=C_TEXT,
                align=PP_ALIGN.LEFT, wrap=True, font=FONT,
                v_anchor=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_label_box(slide, label, l, t, w, h,
                  bg=C_NAVY, text_color=C_WHITE,
                  font_size=Pt(13), bold=True):
    """塗りつぶし矩形＋テキスト"""
    rect = add_rect(slide, l, t, w, h, fill=bg)
    tf = rect.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = text_color
    return rect

def add_multiline_textbox(slide, lines, l, t, w, h,
                          font_size=Pt(12), bold=False, color=C_TEXT,
                          align=PP_ALIGN.LEFT, font=FONT, line_spacing=None):
    """複数行テキスト (リスト of (text, bold, size, color))"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, bold, font_size, color)
        txt, b, sz, cl = item[0], item[1] if len(item)>1 else bold, item[2] if len(item)>2 else font_size, item[3] if len(item)>3 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if line_spacing:
            p.line_spacing = line_spacing
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.name = font
        run.font.size = sz
        run.font.bold = b
        run.font.color.rgb = cl
    return txBox

def rgb_hex(c):
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def gradient_rect(slide, l, t, w, h, color1, color2):
    """グラデーション矩形（縦方向: 上=color1, 下=color2）"""
    rect = add_rect(slide, l, t, w, h, fill=color1)
    sp = rect._element
    spPr = sp.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        spPr.remove(solidFill)
    c1 = rgb_hex(color1)
    c2 = rgb_hex(color2)
    grad_xml = f'''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">
      <a:gsLst>
        <a:gs pos="0">
          <a:srgbClr val="{c1}"/>
        </a:gs>
        <a:gs pos="100000">
          <a:srgbClr val="{c2}"/>
        </a:gs>
      </a:gsLst>
      <a:lin ang="5400000" scaled="0"/>
    </a:gradFill>'''
    grad_elem = parse_xml(grad_xml)
    spPr.insert(list(spPr).index(spPr.find(qn('a:ln'))) if spPr.find(qn('a:ln')) is not None else len(spPr), grad_elem)
    rect.line.fill.background()
    return rect

# =====================================================================
# ページ1: 表紙
# =====================================================================
slide = prs.slides.add_slide(blank_layout)

# 背景（ネイビー→ブルー）
bg = gradient_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_NAVY, C_LTBLUE)

# 上部装飾ライン
add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill=RGBColor(0xFF,0xFF,0xFF))

# "社内資料" ラベル
lbl_w = Inches(2.2)
lbl_h = Inches(0.38)
lbl_l = (SLIDE_W - lbl_w) / 2
lbl_t = Inches(1.4)
rect_lbl = add_rect(slide, lbl_l, lbl_t, lbl_w, lbl_h,
                    fill=None, line=C_WHITE, line_w=Pt(1))
rect_lbl.fill.background()
tf = rect_lbl.text_frame
from pptx.enum.text import MSO_ANCHOR
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "社  内  資  料"
run.font.name = FONT
run.font.size = Pt(12)
run.font.color.rgb = C_WHITE

# メインタイトル
add_textbox(slide, "WC分析ツール",
            Inches(0.5), Inches(2.0), SLIDE_W - Inches(1), Inches(1.3),
            font_size=Pt(54), bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)

# サブタイトル
add_textbox(slide, "重量チェッカーデータ分析システム　紹介資料",
            Inches(0.5), Inches(3.3), SLIDE_W - Inches(1), Inches(0.6),
            font_size=Pt(18), bold=False, color=RGBColor(0xCC, 0xD9, 0xF0),
            align=PP_ALIGN.CENTER)

# 区切り線
add_rect(slide, SLIDE_W/2 - Inches(0.7), Inches(4.05),
         Inches(1.4), Inches(0.04), fill=RGBColor(0xAA, 0xBB, 0xDD))

# キャッチコピー
add_textbox(slide, "重量データから不良品を早期発見し、\n市場流出ゼロを目指す。",
            Inches(0.5), Inches(4.25), SLIDE_W - Inches(1), Inches(1.2),
            font_size=Pt(20), bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER)

# フッター
add_textbox(slide, "2026年6月",
            Inches(0.5), Inches(6.8), SLIDE_W - Inches(1), Inches(0.4),
            font_size=Pt(12), color=RGBColor(0x99, 0xAA, 0xBB),
            align=PP_ALIGN.CENTER)

# =====================================================================
# ページ2: ツール紹介
# =====================================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_WHITE)

# ヘッダーバー
gradient_rect(slide, 0, 0, SLIDE_W, Inches(0.65), C_NAVY, C_BLUE)
add_textbox(slide, "WC分析ツールとは",
            Inches(0.3), Inches(0.08), Inches(8), Inches(0.5),
            font_size=Pt(20), bold=True, color=C_WHITE)

# 概要説明ボックス
add_rect(slide, Inches(0.3), Inches(0.75), SLIDE_W - Inches(0.6), Inches(0.85),
         fill=C_BGWARM, line=C_LTBLUE, line_w=Pt(2))
add_textbox(slide,
            "重量チェッカー（WC）の検査データ（CSV）を読み込み、不良品の自動検出・統計グラフ・エクセルレポートをワンクリックで作成するツールです。\n"
            "OK判定の中に紛れ込んだ異常な重量品も、統計的ばらつき分析によって可視化し、市場への不良品流出を防ぎます。",
            Inches(0.5), Inches(0.82), SLIDE_W - Inches(1.0), Inches(0.7),
            font_size=Pt(12), color=C_TEXT)

# 4機能グリッド
feat_data = [
    ("📂", "CSVファイルを\n読み込むだけ",
     "WCから出力されたCSVを選択するだけで自動分析を開始します。複数ファイルの一括読み込みにも対応。"),
    ("📊", "3種類のグラフを\n自動作成",
     "分布ヒストグラム・OK品時系列・全データ時系列を自動生成。不良発生のタイミングが一目で分かります。"),
    ("🔍", "統計的外れ値を\n自動検出",
     "OK判定品の中から統計的に異常な重量（±3σ超）を自動抽出。WCで弾かれなかった異常品も見逃しません。"),
    ("📝", "エクセルレポートを\n自動作成",
     "分析結果・グラフ・生データをまとめたエクセルファイルをデスクトップに自動保存します。"),
]
col_w = (SLIDE_W - Inches(0.6)) / 2
col_h = Inches(1.55)
for i, (icon, title, desc) in enumerate(feat_data):
    c = i % 2
    r = i // 2
    bx = Inches(0.3) + c * (col_w + Inches(0.1))
    by = Inches(1.75) + r * (col_h + Inches(0.1))
    add_rect(slide, bx, by, col_w, col_h, fill=C_BGSUB,
             line=RGBColor(0xD9, 0xE1, 0xF2), line_w=Pt(1))
    add_textbox(slide, icon, bx + Inches(0.15), by + Inches(0.1),
                Inches(0.5), Inches(0.4), font_size=Pt(18))
    add_textbox(slide, title, bx + Inches(0.65), by + Inches(0.08),
                col_w - Inches(0.8), Inches(0.55),
                font_size=Pt(13), bold=True, color=C_NAVY)
    add_textbox(slide, desc, bx + Inches(0.15), by + Inches(0.65),
                col_w - Inches(0.3), Inches(0.82),
                font_size=Pt(11), color=C_TEXT)

# ランクコード表タイトル
ty = Inches(5.0)
gradient_rect(slide, Inches(0.3), ty, SLIDE_W - Inches(0.6), Inches(0.4), C_NAVY, C_BLUE)
add_textbox(slide, "ランクコードの意味",
            Inches(0.5), ty + Inches(0.04), Inches(6), Inches(0.35),
            font_size=Pt(14), bold=True, color=C_WHITE)

ranks = [
    ("OK",    "正量", "規定重量範囲内の合格品", C_BGSUB, C_TEXT),
    ("軽量",  "軽量", "規定より軽い不合格品。内容物の抜け・欠品が疑われる", C_NG_BG, C_NG_TEXT),
    ("過量",  "過量", "規定より重い不合格品。内容物の過剰投入・異物が疑われる", C_NG_BG, C_NG_TEXT),
    ("２個乗り", "２個乗り", "2個が重なって計量されたと判定されたもの（参考値）", C_WARN_BG, C_WARN_TEXT),
]
row_h = Inches(0.38)
header_y = ty + Inches(0.4)
# ヘッダー行
for j, (lbl, w_frac) in enumerate(zip(["コード","表示名","内容"], [0.15, 0.15, 0.70])):
    col_start = Inches(0.3) + j * (SLIDE_W - Inches(0.6)) * sum([0.15,0.15,0.70][:j])
    col_w2 = (SLIDE_W - Inches(0.6)) * w_frac
    rect = add_rect(slide, col_start, header_y, col_w2, row_h, fill=C_BLUE)
    tf = rect.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = lbl
    run.font.name = FONT; run.font.size = Pt(12); run.font.bold = True
    run.font.color.rgb = C_WHITE

widths = [0.15, 0.15, 0.70]
starts = [Inches(0.3)]
for w in widths[:-1]:
    starts.append(starts[-1] + (SLIDE_W - Inches(0.6)) * w)

for ri, (code, name, desc, bg, tc) in enumerate(ranks):
    ry = header_y + row_h + ri * row_h
    row_bg = bg
    for ci, (txt, cx, cw) in enumerate(zip([code, name, desc], starts, widths)):
        cw2 = (SLIDE_W - Inches(0.6)) * cw
        r2 = add_rect(slide, cx, ry, cw2, row_h, fill=row_bg,
                      line=RGBColor(0xE0, 0xE0, 0xE0), line_w=Pt(0.5))
        tf = r2.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci < 2 else PP_ALIGN.LEFT
        run = p.add_run(); run.text = txt
        run.font.name = FONT; run.font.size = Pt(11); run.font.bold = (ci == 0)
        run.font.color.rgb = tc

# =====================================================================
# ページ3: 使い方
# =====================================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_WHITE)

gradient_rect(slide, 0, 0, SLIDE_W, Inches(0.65), C_NAVY, C_BLUE)
add_textbox(slide, "使い方",
            Inches(0.3), Inches(0.08), Inches(8), Inches(0.5),
            font_size=Pt(20), bold=True, color=C_WHITE)

steps = [
    ("WC分析ツール.exe を起動する",
     "デスクトップにある「WC分析ツール.exe」をダブルクリックして起動します。"),
    ("「CSV選択して解析」ボタンをクリックする",
     "ボタンを押すとファイル選択画面が開きます。重量チェッカーから出力したCSVを選択してください。複数ファイルやフォルダごとの選択も可能です。"),
    ("ロット分割を確認する",
     "データの時間的な切れ目を自動検出し、ロットの区切りを提案します。区切り時間（10分・15分・30分・45分・1時間）をボタンで選んで「この分割でOK」を押します。"),
    ("自動で分析・エクセル作成が完了する",
     "ロット分割を確定すると、自動で統計計算・グラフ作成・エクセル出力が行われます。完了後、デスクトップにファイルが保存されています。"),
    ("エクセルを開いて確認する",
     "「分析レポート」シートに全体統計が、グラフシートには視覚チャートが含まれます。全データ時系列グラフでOK品の中に浮かぶ異常な重量品がないか確認してください。"),
]

step_y = Inches(0.8)
step_h = Inches(1.0)
for i, (title, desc) in enumerate(steps):
    sy = step_y + i * (step_h + Inches(0.08))
    # 丸番号
    circle = add_rect(slide, Inches(0.3), sy, Inches(0.55), Inches(0.55), fill=C_LTBLUE)
    # 円形にする
    sp = circle._element.find(qn('p:spPr'))
    prstGeom = sp.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'ellipse')
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(i+1)
    run.font.name = FONT; run.font.size = Pt(16); run.font.bold = True
    run.font.color.rgb = C_WHITE

    add_textbox(slide, title,
                Inches(1.0), sy, SLIDE_W - Inches(1.3), Inches(0.38),
                font_size=Pt(13), bold=True, color=C_NAVY)
    add_textbox(slide, desc,
                Inches(1.0), sy + Inches(0.38), SLIDE_W - Inches(1.3), Inches(0.58),
                font_size=Pt(11), color=C_TEXT)

# ポイントボックス
note_y = SLIDE_H - Inches(1.3)
add_rect(slide, Inches(0.3), note_y, SLIDE_W - Inches(0.6), Inches(1.1),
         fill=C_WARN_BG, line=RGBColor(0xF9, 0xCA, 0x24), line_w=Pt(1.5))
add_textbox(slide,
            "【ポイント】OK判定でも異常品が混在することがある\n"
            "重量チェッカーは規定範囲内で合否を判定しますが、範囲内でも平均から大きくはずれた品は「外れ値」として本ツールが自動抽出します。\n"
            "時系列チャートの ×印（赤）が外れ値です。同じ時間帯に集中している場合は、その時刻に何らかの異常が発生した可能性があります。",
            Inches(0.5), note_y + Inches(0.08), SLIDE_W - Inches(1.0), Inches(0.95),
            font_size=Pt(11), color=C_WARN_TEXT)

# =====================================================================
# ページ4: 検出事例1・2
# =====================================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_WHITE)

gradient_rect(slide, 0, 0, SLIDE_W, Inches(0.65), C_NAVY, C_BLUE)
add_textbox(slide, "不良品検出事例",
            Inches(0.3), Inches(0.08), Inches(6), Inches(0.5),
            font_size=Pt(20), bold=True, color=C_WHITE)
add_textbox(slide, "WC分析ツールによって未然に防いだ市場流出",
            Inches(6.2), Inches(0.16), Inches(7), Inches(0.35),
            font_size=Pt(12), color=RGBColor(0xAA, 0xBB, 0xCC))

def draw_case(slide, case_num, product, date_str,
              stats_rows, ranks_rows, comment,
              detection, cause, result, result_bg,
              top_y, half_h):
    """事例1枚分を描画"""
    pad = Inches(0.3)
    inner_w = SLIDE_W - pad * 2

    # 事例番号＋品種名
    # 番号丸
    circle = add_rect(slide, pad, top_y + Inches(0.05), Inches(0.58), Inches(0.58), fill=C_NAVY)
    sp = circle._element.find(qn('p:spPr'))
    prstGeom = sp.find(qn('a:prstGeom'))
    if prstGeom is not None: prstGeom.set('prst', 'ellipse')
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = f"事例\n{case_num}"
    run.font.name = FONT; run.font.size = Pt(9); run.font.bold = True
    run.font.color.rgb = C_WHITE

    add_textbox(slide, product, pad + Inches(0.7), top_y + Inches(0.05),
                Inches(6), Inches(0.38),
                font_size=Pt(16), bold=True, color=C_NAVY)
    add_textbox(slide, date_str, pad + Inches(0.7), top_y + Inches(0.43),
                Inches(8), Inches(0.28),
                font_size=Pt(11), color=C_SUBTEXT)

    content_y = top_y + Inches(0.78)
    content_h = half_h - Inches(0.78)
    left_w  = Inches(3.2)
    right_w = inner_w - left_w - Inches(0.15)

    # 統計テーブル (左)
    col_a = Inches(2.0); col_b = Inches(1.2)
    row_h2 = content_h / (len(stats_rows) + 1)
    row_h2 = min(row_h2, Inches(0.32))
    # ヘッダー
    add_rect(slide, pad, content_y, left_w, row_h2, fill=C_BLUE)
    tf2 = slide.shapes[-1].text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run(); run2.text = "基本統計"
    run2.font.name = FONT; run2.font.size = Pt(11); run2.font.bold = True
    run2.font.color.rgb = C_WHITE
    for ri, (key, val, style) in enumerate(stats_rows):
        ry2 = content_y + row_h2 + ri * row_h2
        bg2 = C_NG_BG if style=='ng' else (C_WARN_BG if style=='warn' else (C_BGSUB if ri%2==1 else C_WHITE))
        tc2 = C_NG_TEXT if style=='ng' else (C_WARN_TEXT if style=='warn' else C_TEXT)
        add_rect(slide, pad, ry2, col_a, row_h2, fill=bg2,
                 line=RGBColor(0xE0,0xE0,0xE0), line_w=Pt(0.5))
        r3 = add_rect(slide, pad+col_a, ry2, col_b, row_h2, fill=bg2,
                 line=RGBColor(0xE0,0xE0,0xE0), line_w=Pt(0.5))
        for shape, txt, align in [(slide.shapes[-2], key, PP_ALIGN.LEFT),
                                   (slide.shapes[-1], val, PP_ALIGN.RIGHT)]:
            tf3 = shape.text_frame; tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
            p3 = tf3.paragraphs[0]; p3.alignment = align
            run3 = p3.add_run(); run3.text = txt
            run3.font.name = FONT; run3.font.size = Pt(10)
            run3.font.bold = (style in ('ng','warn'))
            run3.font.color.rgb = tc2

    # ランク＋コメント (右)
    rx = pad + left_w + Inches(0.15)
    add_textbox(slide, "ランクコード集計",
                rx, content_y, right_w, Inches(0.3),
                font_size=Pt(11), bold=True, color=C_NAVY)
    rank_item_h = Inches(0.28)
    for ri2, (rname, rval, rbg, rtc) in enumerate(ranks_rows):
        ry3 = content_y + Inches(0.3) + ri2 * rank_item_h
        add_rect(slide, rx, ry3, right_w, rank_item_h, fill=rbg,
                 line=RGBColor(0xE0,0xE0,0xE0), line_w=Pt(0.3))
        add_textbox(slide, rname, rx+Inches(0.1), ry3, right_w*0.55, rank_item_h,
                    font_size=Pt(10), color=rtc)
        add_textbox(slide, rval, rx, ry3, right_w-Inches(0.1), rank_item_h,
                    font_size=Pt(10), color=rtc, align=PP_ALIGN.RIGHT)

    comment_y = content_y + Inches(0.3) + len(ranks_rows) * rank_item_h + Inches(0.05)
    comment_h = content_h - (comment_y - content_y) - Inches(0.8)
    if comment_h > Inches(0.2):
        add_rect(slide, rx, comment_y, right_w, comment_h, fill=C_BGWARM,
                 line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
        add_textbox(slide, comment, rx+Inches(0.1), comment_y+Inches(0.05),
                    right_w-Inches(0.2), comment_h-Inches(0.1),
                    font_size=Pt(10), color=C_TEXT)

    # 検出/原因/結果 カード
    card_y = top_y + half_h - Inches(0.75)
    card_w = inner_w / 3 - Inches(0.08)
    cards = [
        ("検出内容", detection, C_WARN_BG, C_WARN_TEXT),
        ("判明した異常", cause, C_NG_BG, C_NG_TEXT),
        ("防いだ結果", result, result_bg, C_OK_TEXT if result_bg==C_OK_BG else C_NAVY),
    ]
    for ci2, (clabel, cval, cbg, ctc) in enumerate(cards):
        cx2 = pad + ci2 * (card_w + Inches(0.12))
        add_rect(slide, cx2, card_y, card_w, Inches(0.72), fill=cbg,
                 line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
        add_textbox(slide, clabel, cx2+Inches(0.1), card_y+Inches(0.04),
                    card_w-Inches(0.2), Inches(0.22),
                    font_size=Pt(9), color=C_SUBTEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, cval, cx2+Inches(0.1), card_y+Inches(0.26),
                    card_w-Inches(0.2), Inches(0.42),
                    font_size=Pt(11), bold=True, color=ctc, align=PP_ALIGN.CENTER)

# ----- 事例1 -----
draw_case(slide,
    case_num="1",
    product="コンタック 36錠",
    date_str="製造日：2025年12月22〜23日　／　品種番号11　ロット1",
    stats_rows=[
        ("全数",          "19,715 個",  ""),
        ("OK数",          "19,682 個",  ""),
        ("NG数",          "33 個",      "ng"),
        ("不良率",        "0.17 %",     "warn"),
        ("平均重量",      "34.599 g",   ""),
        ("標準偏差",      "0.202 g",    ""),
        ("外れ値件数",    "41 件",      "warn"),
    ],
    ranks_rows=[
        ("OK（正量）",  "19,682件 99.83%", C_BGSUB,    C_TEXT),
        ("２個乗り",    "21件 0.11%",      C_WARN_BG,  C_WARN_TEXT),
        ("過量",        "7件 0.04%",       C_NG_BG,    C_NG_TEXT),
        ("軽量",        "5件 0.03%",       C_NG_BG,    C_NG_TEXT),
    ],
    comment="全データ時系列で、OK品（約34.6g）の中に10〜20g台という極端に軽い点が複数散見された。WCはOK判定していたが重量の乖離を検出。",
    detection="OK品の中に\n極端に軽い個体を検出",
    cause="PTPシート\n1枚抜け",
    result="不良品の市場流出を防止\n速やかに工程対応",
    result_bg=C_OK_BG,
    top_y=Inches(0.68),
    half_h=Inches(3.35),
)

# 区切り線
add_rect(slide, Inches(0.3), Inches(4.03), SLIDE_W - Inches(0.6), Inches(0.03),
         fill=RGBColor(0xE0,0xD8,0xCC))

# ----- 事例2 -----
draw_case(slide,
    case_num="2",
    product="防風通聖散料 360錠",
    date_str="製造日：2026年3月30日　／　品種番号20　ロット1",
    stats_rows=[
        ("全数",          "2,172 個",   ""),
        ("OK数",          "2,160 個",   ""),
        ("NG数",          "12 個",      "ng"),
        ("不良率",        "0.55 %",     "warn"),
        ("平均重量",      "383.023 g",  ""),
        ("標準偏差",      "0.510 g",    ""),
        ("外れ値件数",    "6 件",       "warn"),
    ],
    ranks_rows=[
        ("OK（正量）",  "2,160件 99.45%", C_BGSUB,    C_TEXT),
        ("２個乗り",    "10件 0.46%",     C_WARN_BG,  C_WARN_TEXT),
        ("軽量",        "2件 0.09%",      C_NG_BG,    C_NG_TEXT),
    ],
    comment="全データ時系列で、383g付近のOK品の中にほぼ0g近くの測定値が製造開始直後に2点検出された。添付文書の入れ忘れによる重量差を可視化。",
    detection="OK品の中に\n極端に軽い個体を検出",
    cause="添付文書\n1枚抜け",
    result="不良品の市場流出を防止\n工程の早期対応を実施",
    result_bg=C_OK_BG,
    top_y=Inches(4.06),
    half_h=Inches(3.35),
)

# =====================================================================
# ページ5: 検出事例3・4
# =====================================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_WHITE)

gradient_rect(slide, 0, 0, SLIDE_W, Inches(0.65), C_NAVY, C_BLUE)
add_textbox(slide, "不良品検出事例",
            Inches(0.3), Inches(0.08), Inches(6), Inches(0.5),
            font_size=Pt(20), bold=True, color=C_WHITE)
add_textbox(slide, "WC分析ツールによって未然に防いだ市場流出",
            Inches(6.2), Inches(0.16), Inches(7), Inches(0.35),
            font_size=Pt(12), color=RGBColor(0xAA, 0xBB, 0xCC))

draw_case(slide,
    case_num="3",
    product="第一三共胃腸薬プラス細粒 12包",
    date_str="製造日：2026年6月22日　／　品種番号1　ロット1",
    stats_rows=[
        ("全数",          "22,827 個",  ""),
        ("OK数",          "22,804 個",  ""),
        ("NG数",          "23 個",      "ng"),
        ("不良率",        "0.10 %",     ""),
        ("平均重量",      "36.976 g",   ""),
        ("標準偏差",      "0.153 g",    ""),
        ("外れ値件数",    "124 件",     "warn"),
    ],
    ranks_rows=[
        ("OK（正量）",  "22,804件 99.90%", C_BGSUB,    C_TEXT),
        ("軽量",        "16件 0.07%",      C_NG_BG,    C_NG_TEXT),
        ("過量",        "5件 0.02%",       C_NG_BG,    C_NG_TEXT),
        ("２個乗り",    "2件 0.01%",       C_WARN_BG,  C_WARN_TEXT),
    ],
    comment="全データ時系列で、OK品（約37g）の中に約40gという突出した重量の点が散発的に検出された。細粒1包が余分に挿入されていた重量増加を可視化。",
    detection="OK品の中に\n極端に重い個体を検出",
    cause="細粒が\n1包多く挿入",
    result="不良品の市場流出を防止\n充填工程を即時確認",
    result_bg=C_OK_BG,
    top_y=Inches(0.68),
    half_h=Inches(3.35),
)

add_rect(slide, Inches(0.3), Inches(4.03), SLIDE_W - Inches(0.6), Inches(0.03),
         fill=RGBColor(0xE0,0xD8,0xCC))

draw_case(slide,
    case_num="4",
    product="ルプセ",
    date_str="製造日：2026年3月27日　／　品種番号7　ロット1",
    stats_rows=[
        ("全数",          "3,641 個",   ""),
        ("OK数",          "3,638 個",   ""),
        ("NG数",          "3 個",       ""),
        ("不良率",        "0.08 %",     ""),
        ("平均重量",      "116.746 g",  ""),
        ("標準偏差",      "0.230 g",    ""),
        ("外れ値件数",    "98 件",      "ng"),
    ],
    ranks_rows=[
        ("OK（正量）",  "3,638件 99.92%", C_BGSUB,    C_TEXT),
        ("２個乗り",    "2件 0.05%",      C_WARN_BG,  C_WARN_TEXT),
        ("軽量",        "1件 0.03%",      C_NG_BG,    C_NG_TEXT),
    ],
    comment="NG数はわずか3件だったが、OK判定品の中に外れ値が98件検出された。時系列チャートで×印が特定の時間帯に集中しており、添付文書が2枚入りと判明。",
    detection="OK品の中から\n外れ値を98件検出",
    cause="添付文書が\n2枚入り",
    result="PQロットのため市場流出なし\n異常の記録・工程改善に活用",
    result_bg=C_BGSUB,
    top_y=Inches(4.06),
    half_h=Inches(3.35),
)

# =====================================================================
# ページ6: まとめ
# =====================================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=C_WHITE)

gradient_rect(slide, 0, 0, SLIDE_W, Inches(0.65), C_NAVY, C_BLUE)
add_textbox(slide, "まとめ",
            Inches(0.3), Inches(0.08), Inches(8), Inches(0.5),
            font_size=Pt(20), bold=True, color=C_WHITE)

# まとめテーブル
headers = ["事例", "品種", "製造日", "検出内容", "判明した異常", "結果"]
col_ws  = [0.06, 0.22, 0.12, 0.26, 0.20, 0.14]
rows_data = [
    ("1", "コンタック36錠",          "2025/12/22", "OK品の中の極端な軽量品", "PTPシート1枚抜け",   "市場流出防止",       "ok"),
    ("2", "防風通聖散料360錠",        "2026/3/30",  "OK品の中の極端な軽量品", "添付文書1枚抜け",    "市場流出防止",       "ok"),
    ("3", "第一三共胃腸薬プラス細粒12包", "2026/6/22","OK品の中の極端な重量品", "細粒1包多く挿入",   "市場流出防止",       "ok"),
    ("4", "ルプセ",                   "2026/3/27",  "OK品中の外れ値98件",     "添付文書2枚入り",    "PQ品、記録に活用",   "blue"),
]
total_w = SLIDE_W - Inches(0.6)
starts2 = [Inches(0.3)]
for cw3 in col_ws[:-1]:
    starts2.append(starts2[-1] + total_w * cw3)

header_y2 = Inches(0.75)
row_h3 = Inches(0.42)

for ci3, (hdr, cw3, cx3) in enumerate(zip(headers, col_ws, starts2)):
    cw3v = total_w * cw3
    add_rect(slide, cx3, header_y2, cw3v, row_h3, fill=C_BLUE)
    tf4 = slide.shapes[-1].text_frame
    tf4.vertical_anchor = MSO_ANCHOR.MIDDLE
    p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run(); run4.text = hdr
    run4.font.name = FONT; run4.font.size = Pt(12); run4.font.bold = True
    run4.font.color.rgb = C_WHITE

for ri3, row3 in enumerate(rows_data):
    ry4 = header_y2 + row_h3 + ri3 * row_h3
    row_bg2 = C_BGSUB if ri3 % 2 == 1 else C_WHITE
    for ci4, (val4, cw4, cx4) in enumerate(zip(row3[:-1], col_ws, starts2)):
        cw4v = total_w * cw4
        is_result = (ci4 == 5)
        style4 = row3[-1]
        cell_bg = (C_OK_BG if style4=='ok' else C_BGSUB) if is_result else row_bg2
        cell_tc = (C_OK_TEXT if style4=='ok' else C_NAVY) if is_result else C_TEXT
        add_rect(slide, cx4, ry4, cw4v, row_h3, fill=cell_bg,
                 line=RGBColor(0xE0,0xE0,0xE0), line_w=Pt(0.5))
        tf5 = slide.shapes[-1].text_frame
        tf5.vertical_anchor = MSO_ANCHOR.MIDDLE
        p5 = tf5.paragraphs[0]
        p5.alignment = PP_ALIGN.CENTER if ci4 in (0, 2) else PP_ALIGN.LEFT
        run5 = p5.add_run(); run5.text = val4
        run5.font.name = FONT; run5.font.size = Pt(11)
        run5.font.bold = is_result
        run5.font.color.rgb = cell_tc

# 実績ハイライト
hl_y = Inches(2.6)
highlight = gradient_rect(slide, Inches(0.3), hl_y, SLIDE_W - Inches(0.6), Inches(1.2), C_NAVY, C_BLUE)
add_textbox(slide, "これまでの実績",
            Inches(0.3), hl_y + Inches(0.05), SLIDE_W - Inches(0.6), Inches(0.3),
            font_size=Pt(12), color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(slide, "4件の不良品流出を防止",
            Inches(0.3), hl_y + Inches(0.35), SLIDE_W - Inches(0.6), Inches(0.55),
            font_size=Pt(32), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "いずれも重量チェッカーのOK判定をすり抜けていた異常品",
            Inches(0.3), hl_y + Inches(0.9), SLIDE_W - Inches(0.6), Inches(0.32),
            font_size=Pt(12), color=RGBColor(0xCC, 0xD9, 0xF0), align=PP_ALIGN.CENTER)

# まとめ文
sum_y = hl_y + Inches(1.35)
add_rect(slide, Inches(0.3), sum_y, SLIDE_W - Inches(0.6), Inches(2.1),
         fill=C_BGWARM, line=C_LTBLUE, line_w=Pt(2))
add_textbox(slide,
            "WC分析ツールは、重量チェッカーの判定だけでは気づけなかった異常を、統計的なアプローチで可視化します。\n\n"
            "製造ロットごとに全数データを確認することで、発生した異常の「時間帯・傾向・件数」を把握し、工程の早期改善につなげることができます。\n\n"
            "日々の製造データを分析習慣に取り入れることで、品質のさらなる安定と、市場流出ゼロの継続を目指しましょう。",
            Inches(0.55), sum_y + Inches(0.15), SLIDE_W - Inches(1.1), Inches(1.8),
            font_size=Pt(13), color=C_TEXT)

# =====================================================================
# 保存
# =====================================================================
out_path = r"c:\Users\k-nakamori\Desktop\WC分析ツール_紹介資料.pptx"
prs.save(out_path)
print(f"保存完了: {out_path}")
